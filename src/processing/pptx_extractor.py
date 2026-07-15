"""
src/processing/pptx_extractor.py
────────────────────────────────────────────────────────────
Extract meaningful keyword candidates from a .pptx file.

Used by the Streamlit keyword_manager.py — the user uploads a deck,
this module returns a list of cleaned tokens which are then previewed
and confirmed before being added to the keywords table.
────────────────────────────────────────────────────────────
"""

import re
import logging
from typing import List

logger = logging.getLogger(__name__)

# Common English stop-words to filter out before returning tokens
_STOPWORDS = frozenset({
    "a", "an", "the", "and", "or", "but", "in", "on", "at", "to", "for",
    "of", "with", "by", "from", "is", "are", "was", "were", "be", "been",
    "being", "have", "has", "had", "do", "does", "did", "will", "would",
    "could", "should", "may", "might", "shall", "can", "not", "no", "nor",
    "so", "yet", "both", "either", "neither", "as", "if", "then", "than",
    "that", "this", "these", "those", "it", "its", "we", "our", "you",
    "your", "he", "his", "she", "her", "they", "their", "all", "any",
    "each", "every", "more", "most", "other", "some", "such", "into",
    "through", "during", "before", "after", "above", "below", "between",
    "out", "off", "over", "under", "again", "further", "once",
})

_MIN_TOKEN_LEN = 3


def extract_keywords_from_pptx(filepath: str) -> List[str]:
    """
    Parse all text frames in every slide and return a de-duplicated,
    cleaned list of candidate keywords.

    Parameters
    ----------
    filepath : str
        Absolute or relative path to the .pptx file.

    Returns
    -------
    List[str]
        Sorted list of unique keyword candidates (lowercased).
    """
    try:
        from pptx import Presentation  # lazy import
    except ImportError:
        logger.error("python-pptx not installed — run: pip install python-pptx")
        return []

    candidates: set = set()

    try:
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        raw = run.text.strip()
                        if not raw:
                            continue
                        # Tokenise: split on whitespace + punctuation
                        tokens = re.split(r"[\s\-–—,;:!?.\"'()\[\]{}<>]+", raw)
                        for tok in tokens:
                            tok = tok.lower().strip()
                            if (
                                len(tok) >= _MIN_TOKEN_LEN
                                and tok not in _STOPWORDS
                                and not tok.isdigit()
                                and re.search(r"[a-z]", tok)   # at least one letter
                            ):
                                candidates.add(tok)

    except Exception as exc:
        logger.error("pptx extraction failed: %s", exc)

    result = sorted(candidates)
    logger.info("Extracted %d keyword candidates from %s", len(result), filepath)
    return result


def extract_keywords_from_bytes(pptx_bytes: bytes) -> List[str]:
    """
    Same as extract_keywords_from_pptx but accepts raw bytes (for Streamlit
    st.file_uploader which returns a BytesIO-compatible object).
    """
    import io
    try:
        from pptx import Presentation
    except ImportError:
        logger.error("python-pptx not installed")
        return []

    candidates: set = set()
    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        raw = run.text.strip()
                        if not raw:
                            continue
                        tokens = re.split(r"[\s\-–—,;:!?.\"'()\[\]{}<>]+", raw)
                        for tok in tokens:
                            tok = tok.lower().strip()
                            if (
                                len(tok) >= _MIN_TOKEN_LEN
                                and tok not in _STOPWORDS
                                and not tok.isdigit()
                                and re.search(r"[a-z]", tok)
                            ):
                                candidates.add(tok)
    except Exception as exc:
        logger.error("pptx byte extraction failed: %s", exc)

    return sorted(candidates)
